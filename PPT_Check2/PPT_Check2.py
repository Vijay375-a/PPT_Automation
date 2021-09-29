from pptx import Presentation
import re
import sys
from datetime import datetime

iter = sys.argv[1]
iter1 = sys.argv[2]
start = datetime.now()
print("--------------------------------------------------------------------------------------------------------")
print("Document Name:", iter)
print("CheckList Rule - 2: Validation of font name")
print("Document Review Start Time:", start, "HH:MM:SS")
print("--------------------------------------------------------------------------------------------------------")
print("\n")
prs = Presentation(iter)
# text_runs will be populated with a list of strings,
# one for each text run in presentation
text_font = []
text_font_name = []  # Common font names
cnt = 0
cnt1 = 0
text_font_found = []
text_found = []
text_modified = []
for slide in prs.slides:  # Reading titles related data and paragraph data
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    # print(run.text)
                    if run.font.name is not None:
                        text_font.append(run.font.name)

                        if run.font.name == iter1:
                            text_font_found.append(run.font.name)
                            cnt = cnt + 1
# print(text_font_found)
for i in text_font:
    if i not in text_font_name:
        text_font_name.append(i)

for i in text_font_found:
    if i not in text_found:
        text_found.append(i)

text_modified = [x for x in text_font_name if x not in text_found]

if cnt >= 1:
    print("Status:Pass")
    for name1 in text_modified:
        print("Other font type used:", name1)


else:
    print("Status:Fail")

    for name in text_font_name:
        print("This are the font names used in PPT file:", name)
end = datetime.now()
print("\nDocument Review End Time:", end)
print("\nTime taken For Document Review:", end - start, "HH:MM:SS")
