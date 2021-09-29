from pptx import Presentation
import re
import sys
from datetime import datetime

iter = sys.argv[1]
start = datetime.now()
print("--------------------------------------------------------------------------------------------------------")
print("Document Name:", iter)
print("CheckList Rule - 1: Check  Date Format (Supported date formats are yyyy/mm/dd or yyyy-mm-dd)")
print("Document Review Start Time:", start, "HH:MM:SS")
print("--------------------------------------------------------------------------------------------------------")
print("\n")
prs = Presentation(iter)
text_title = []
text_runs = []
l = []
l3 = []
l4 = []
cnt = 0
cnt1 = 0
cnt2 = 0
cnt3 = 0
for slide in prs.slides:  # Reading paragraph data for PPT files(Regex : \d{4}[/-]\d{2}[/-]\d{2})
    try:
        title = slide.shapes.title.text
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
                        m1 = str(run.text)
                        m2 = m1.rstrip('\r\x07')
                        m = m2.rstrip('\t\r')
                        if len(m) != 0:
                            match1 = re.search(r'\d{4}[/-]\d{2}[/-]\d{2}', m)
                            if match1 is not None:
                                l.append(match1.group())
                                # print(l)
                                # print("Title name:", title)

                                try:
                                    cnt = cnt + 1
                                    date1 = datetime.strptime(match1.group(), '%Y-%m-%d').date()
                                    # print("Date found:", date1)
                                except ValueError:
                                    cnt = cnt + 1
                                    date1 = datetime.strptime(match1.group(), '%Y/%m/%d').date()

    except AttributeError:
        pass
# print("Date found:", l[-1])

for slide in prs.slides:
    try:
        title = slide.shapes.title.text  # Reading tables data for PPT files(Regex : \d{4}[/-]\d{2}[/-]\d{2})
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                for cell in table.iter_cells():
                    text_title.append(cell.text)
                    m1 = str(cell.text)
                    m2 = m1.rstrip('\r\x07')
                    m = m2.rstrip('\t\r')
                    if len(m) != 0:
                        match1 = re.search(r'\d{4}[/-]\d{2}[/-]\d{2}', m)
                        if match1 is not None:
                            l.append(match1.group())
                            # print("Title name:", title)

                            try:
                                cnt1 = cnt1 + 1
                                date1 = datetime.strptime(match1.group(), '%Y-%m-%d').date()
                                # print("Date found:", date1)
                            except ValueError:
                                cnt1 = cnt1 + 1
                                date1 = datetime.strptime(match1.group(), '%Y/%m/%d').date()
                            # print("Date found:", l[-1])
    except AttributeError:
        pass
for slide in prs.slides:  # Reading paragraph data for PPT files(Regex : \d{4}[/-]\S{3}[/-]\d{2})
    try:
        title = slide.shapes.title.text
        # Reading titles related data
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
                        m1 = str(run.text)
                        m2 = m1.rstrip('\r\x07')
                        m = m2.rstrip('\t\r')
                        if len(m) != 0:
                            match1 = re.search(r'\d{4}[/-]\S{3}[/-]\d{2}', m)
                            if match1 is not None:
                                l.append(match1.group())
                                cnt = cnt +1

    except AttributeError:
        pass
for slide in prs.slides:  # Reading tables data for PPT files(Regex : \d{4}[/-]\S{3}[/-]\d{2})
    try:
        title = slide.shapes.title.text
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                for cell in table.iter_cells():
                    text_title.append(cell.text)
                    m1 = str(cell.text)
                    m2 = m1.rstrip('\r\x07')
                    m = m2.rstrip('\t\r')
                    if len(m) != 0:
                        match1 = re.search(r'\d{4}[/-]\S{3}[/-]\d{2}', m)
                        if match1 is not None:
                            l.append(match1.group())
                            # print("Title name:", title)
                            cnt1 = cnt1 + 1

    except AttributeError:
        pass
for slide in prs.slides:  # Reading tables data for PPT files(Regex : \d{4}[.]\S{3}[.]\d{2})
    try:
        title = slide.shapes.title.text
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                for cell in table.iter_cells():
                    text_title.append(cell.text)
                    m1 = str(cell.text)
                    m2 = m1.rstrip('\r\x07')
                    m = m2.rstrip('\t\r')
                    if len(m) != 0:
                        match1 = re.search(r'\d{4}[.]\S{3}[.]\d{2}', m)
                        if match1 is not None:
                            l3.append(match1.group())
                            # print(l3)
                            print("Title name:", title)
                            cnt2 = cnt2 + 1
    except  AttributeError:
        pass
for slide in prs.slides:  # Reading paragraph data for PPT files(Regex : \d{4}[.]\S{3}[.]\d{2})
    try:
        title = slide.shapes.title.text
        # Reading titles related data
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
                        m1 = str(run.text)
                        m2 = m1.rstrip('\r\x07')
                        m = m2.rstrip('\t\r')
                        if len(m) != 0:
                            match1 = re.search(r'\d{4}[.]\S{3}[.]\d{2}', m)
                            if match1 is not None:
                                l3.append(match1.group())
                                # print(l)
                                print("Title name:", title)
                                cnt3 = cnt3 + 1
    except AttributeError:
        pass

for slide in prs.slides:  # Reading paragraph data for PPT files(Regex : \d{4}[/-]\d{2}[/-]\d{2})
    try:
        title = slide.shapes.title.text
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
                        m1 = str(run.text)
                        m2 = m1.rstrip('\r\x07')
                        m = m2.rstrip('\t\r')
                        if len(m) != 0:
                            match1 = re.search(r'\d{4}[.]\d{2}[.]\d{2}', m)
                            if match1 is not None:
                                l.append(match1.group())
                                # print(l)
                                print("Title name:", title)
                                cnt2 = cnt2 + 1
    except AttributeError:
        pass
for slide in prs.slides:
    try:
        title = slide.shapes.title.text  # Reading tables data for PPT files(Regex : \d{4}[/-]\d{2}[/-]\d{2})
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                for cell in table.iter_cells():
                    text_title.append(cell.text)
                    m1 = str(cell.text)
                    m2 = m1.rstrip('\r\x07')
                    m = m2.rstrip('\t\r')
                    if len(m) != 0:
                        match1 = re.search(r'\d{4}[.]\d{2}[.]\d{2}', m)
                        if match1 is not None:
                            l.append(match1.group())
                            print("Title name:", title)
                            cnt3 = cnt3 + 1
    except AttributeError:
        pass
for slide in prs.slides:  # Reading paragraph data for PPT files(Regex :\d{8})
    try:
        title = slide.shapes.title.text
        # Reading titles related data
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
                        m1 = str(run.text)
                        m2 = m1.rstrip('\r\x07')
                        m = m2.rstrip('\t\r')
                        if len(m) != 0:
                            match1 = re.search(r'\d{8}', m)
                            if match1 is not None:
                                l4.append(match1.group())
                                # print(l4)
                                print("Title name:", title)
                                cnt3 = cnt3 + 1

    except AttributeError:
        pass
for slide in prs.slides:  # Reading paragraph data for PPT files(Regex :\d{8})
    try:
        title = slide.shapes.title.text
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                for cell in table.iter_cells():
                    text_title.append(cell.text)
                    m1 = str(cell.text)
                    m2 = m1.rstrip('\r\x07')
                    m = m2.rstrip('\t\r')
                    if len(m) != 0:
                        match1 = re.search(r'\d{8}', m)
                        if match1 is not None:
                            l3.append(match1.group())
                            # print(l3)
                            print("Title name:", title)
                            cnt2 = cnt2 + 1
    except AttributeError:
        pass
if cnt2 >= 1 or cnt3 >= 1:
    print("Status:Fail")
elif cnt >= 1 and cnt1 >= 1:
    print("Status:Pass")

end = datetime.now()
print("\nDocument Review End Time:", end)
print("\nTime taken For Document Review:", end - start, "HH:MM:SS")
