from pptx import Presentation
import re
import sys
from datetime import datetime

iter = sys.argv[1]
iter1 = sys.argv[2]
start = datetime.now()
print("--------------------------------------------------------------------------------------------------------")
print("Document Name:", iter)
print("CheckList Rule - 3: Check for T.B.D")
print("Document Review Start Time:", start, "HH:MM:SS")
print("--------------------------------------------------------------------------------------------------------")
print("\n")

l = []
cnt = 0
cnt1 = 0
text_runs = []
prs = Presentation(iter)
for slide in prs.slides:
    try:
        title = slide.shapes.title.text
        # Reading titles related data
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
                        m1 = str(run.text)
                        m2 = m1.rstrip('\r\x07')
                        m = m2.rstrip('\t\r')
                        if len(m) != 0:
                            match1 = re.search(iter1, m)

                            # print(match1)
                            if match1 is not None:
                                l.append(match1.group())

                                print("Title name:", title)
                                try:
                                    cnt = cnt + 1
                                except:

                                    pass


    except AttributeError:
        print()
for slide in prs.slides:  # Code for Searching TBD in PPT file.
    try:
        title = slide.shapes.title.text
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                for cell in table.iter_cells():
                    m1 = str(cell.text)
                    m2 = m1.rstrip('\r\x07')
                    m = m2.rstrip('\t\r')
                    if len(m) != 0:
                        match1 = re.search(iter1, m)

                        # print(match1)
                        if match1 is not None:
                            l.append(match1.group())

                            print("Title name:", title)
                            try:
                                cnt = cnt + 1
                            except:

                                pass
    except AttributeError:
        print()

if cnt >= 1 or cnt1 >= 1:
    print("Status:Fail")
else:
    print("Status:Pass")
    print("No TBD found in PPT files")
end = datetime.now()
print("\nDocument Review End Time:", end)
print("\nTime taken For Document Review:", end - start, "HH:MM:SS")
